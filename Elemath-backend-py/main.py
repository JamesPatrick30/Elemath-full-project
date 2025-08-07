from fastapi import FastAPI, File, UploadFile, Form
from fastapi.responses import JSONResponse
import os
import pdfplumber
from docx import Document
from pptx import Presentation
from openai import OpenAI
import shutil

app = FastAPI()

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Initialize OpenRouter Client
client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key="YOUR_OPENROUTER_API_KEY",  # Replace this
)

# Utility functions to extract text
def extract_text_from_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text(file_path, filename):
    ext = filename.split('.')[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_path)
    elif ext == "docx":
        return extract_text_from_docx(file_path)
    elif ext == "pptx":
        return extract_text_from_pptx(file_path)
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    else:
        return ""

# FastAPI Route
@app.post("/upload/")
async def upload_file(file: UploadFile = File(...), prompt: str = Form("Summarize this text.")):
    filename = file.filename
    file_path = os.path.join(UPLOAD_FOLDER, filename)

    with open(file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    raw_text = extract_text(file_path, filename)

    if not raw_text.strip():
        return JSONResponse({"error": "Failed to extract text."}, status_code=400)

    # Send to OpenRouter
    try:
        completion = client.chat.completions.create(
            model="mistralai/mistral-nemo:free",
            messages=[
                {"role": "system", "content": "You are a helpful assistant."},
                {"role": "user", "content": f"{prompt}\n\n{raw_text}"}
            ],
        )
        ai_response = completion.choices[0].message.content
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

    return {"result": ai_response}

